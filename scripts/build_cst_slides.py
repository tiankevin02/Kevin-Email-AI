#!/usr/bin/env python3
"""
CST トラクトグラフィー解析手順スライド生成スクリプト
SPM/MRtrix3 を用いた皮質脊髄路(CST)トラクトグラフィーのやり方まとめPDFに基づく。
スクリーンショットは後ほど差し込むため、プレースホルダ枠を配置する。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- デザイン定数 ----
NAVY   = RGBColor(0x16, 0x2A, 0x47)   # 濃紺(タイトル/帯)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)   # アクセント青
CYAN   = RGBColor(0x35, 0xC7, 0xEC)   # シアン(表紙の色)
TEAL   = RGBColor(0x00, 0x9E, 0xB8)
LIGHT  = RGBColor(0xF2, 0xF6, 0xFA)   # 薄い背景
GRAY   = RGBColor(0x5A, 0x6B, 0x7B)
DARK   = RGBColor(0x1B, 0x24, 0x2E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG = RGBColor(0x1E, 0x29, 0x39)   # コードブロック背景
CODEFG = RGBColor(0xCB, 0xE6, 0xFF)
PH_BG  = RGBColor(0xE6, 0xEC, 0xF2)   # プレースホルダ背景
PH_LN  = RGBColor(0xB4, 0xC2, 0xD0)

JP   = "Yu Gothic"          # 日本語フォント
JPB  = "Yu Gothic UI Semibold"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def set_font(run, size=None, bold=None, color=None, name=JP, italic=None):
    f = run.font
    if size  is not None: f.size = Pt(size)
    if bold  is not None: f.bold = bold
    if italic is not None: f.italic = italic
    if color is not None: f.color.rgb = color
    f.name = name
    # 東アジア言語フォントも明示
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp


def textbox(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of dicts -> {text,size,bold,color,name,space_after,bullet,italic}"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get('align', align)
        if ln.get('space_after') is not None: p.space_after = Pt(ln['space_after'])
        if ln.get('space_before') is not None: p.space_before = Pt(ln['space_before'])
        if ln.get('line') is not None: p.line_spacing = ln['line']
        run = p.add_run()
        run.text = ln['text']
        set_font(run, ln.get('size', 18), ln.get('bold'), ln.get('color', DARK),
                 ln.get('name', JP), ln.get('italic'))
    return tb


def bg(s, color=WHITE):
    rect(s, 0, 0, SW, SH, fill=color)


def header(s, step_label, title, idx=None):
    """共通ヘッダー(左帯+ステップラベル+タイトル)"""
    rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    rect(s, 0, 0, Inches(0.18), SH, fill=CYAN)
    if step_label:
        chip = rect(s, Inches(0.55), Inches(0.30), Inches(2.0), Inches(0.55),
                    fill=CYAN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        chip.adjustments[0] = 0.5
        textbox(s, Inches(0.55), Inches(0.30), Inches(2.0), Inches(0.55),
                [{'text': step_label, 'size': 16, 'bold': True, 'color': NAVY}],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = Inches(2.75)
    else:
        tx = Inches(0.55)
    textbox(s, tx, 0, SW - tx - Inches(0.5), Inches(1.15),
            [{'text': title, 'size': 28, 'bold': True, 'color': WHITE}],
            anchor=MSO_ANCHOR.MIDDLE)
    if idx is not None:
        textbox(s, SW - Inches(1.2), Inches(0.32), Inches(0.9), Inches(0.5),
                [{'text': idx, 'size': 12, 'color': RGBColor(0xAF,0xC6,0xDB)}],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def placeholder(s, x, y, w, h, label="スクリーンショット挿入予定"):
    box = rect(s, x, y, w, h, fill=PH_BG, line=PH_LN, line_w=Pt(1.25),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.03
    textbox(s, x, y, w, h,
            [{'text': '🖼', 'size': 30, 'color': PH_LN, 'align': PP_ALIGN.CENTER,
              'space_after': 4},
             {'text': label, 'size': 13, 'bold': True, 'color': GRAY, 'align': PP_ALIGN.CENTER}],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return box


def add_image_fit(s, path, x, y, w, h):
    """アスペクト比を保って枠内に収め、中央配置する"""
    from PIL import Image as PImage
    iw, ih = PImage.open(path).size
    ar = iw / ih
    br = int(w) / int(h)
    if ar > br:
        nw = int(w); nh = int(int(w) / ar)
    else:
        nh = int(h); nw = int(int(h) * ar)
    nx = int(x) + (int(w) - nw) // 2
    ny = int(y) + (int(h) - nh) // 2
    return s.shapes.add_picture(path, nx, ny, nw, nh)


def codebox(s, x, y, w, h, cmds, title=None):
    box = rect(s, x, y, w, h, fill=CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box.adjustments[0] = 0.04
    # 信号機ドット
    for i, c in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
        d = rect(s, x + Inches(0.18) + Inches(0.22)*i, y + Inches(0.16),
                 Inches(0.12), Inches(0.12), fill=c, shape=MSO_SHAPE.OVAL)
    tb = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.42), w - Inches(0.4), h - Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    for i, c in enumerate(cmds):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.1; p.space_after = Pt(3)
        prompt = p.add_run(); prompt.text = "$ "
        set_font(prompt, 13, True, CYAN, MONO)
        run = p.add_run(); run.text = c
        set_font(run, 13, False, CODEFG, MONO)
    return box


# ============================================================
# 1. 表紙
# ============================================================
s = slide(); bg(s, NAVY)
rect(s, 0, 0, SW, SH, fill=NAVY)
# 斜めのシアン帯
rect(s, 0, Inches(4.7), SW, Inches(0.12), fill=CYAN)
rect(s, 0, Inches(5.0), Inches(6.0), Inches(0.06), fill=TEAL)
textbox(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(2.6),
        [{'text': 'CST トラクトグラフィー', 'size': 54, 'bold': True, 'color': WHITE, 'space_after': 6},
         {'text': '解析手順まとめ', 'size': 40, 'bold': True, 'color': CYAN}])
textbox(s, Inches(0.95), Inches(5.15), Inches(11), Inches(1.3),
        [{'text': '皮質脊髄路 (Corticospinal Tract) を MRtrix3 で抽出する', 'size': 20, 'color': RGBColor(0xCF,0xDD,0xEA), 'space_after': 4},
         {'text': 'DWI → 全脳トラクト生成 → ROI による絞り込み → 可視化', 'size': 15, 'color': RGBColor(0x9D,0xB4,0xC8)}])
textbox(s, Inches(0.95), Inches(6.6), Inches(11), Inches(0.5),
        [{'text': '2026-06-16  /  Tutorial: MRI/DWI/CST', 'size': 13, 'color': RGBColor(0x7F,0x96,0xAB)}])

# ============================================================
# 2. このスライドの目的 / CSTとは
# ============================================================
s = slide(); bg(s)
header(s, None, "はじめに — CST と本資料のゴール", "2")
textbox(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(5.3),
        [{'text': '皮質脊髄路 (CST) とは', 'size': 22, 'bold': True, 'color': BLUE, 'space_after': 8},
         {'text': '一次運動野 (M1) から脊髄へ運動指令を伝える主要な下行性伝導路。随意運動の中枢で、解剖学的に明確な走行をもつ。', 'size': 16, 'color': DARK, 'space_after': 16, 'line': 1.25},
         {'text': '本資料のゴール', 'size': 22, 'bold': True, 'color': BLUE, 'space_after': 8},
         {'text': '拡散強調画像 (DWI) から CST のみを抽出し、3D トラクトグラフィーとして可視化する一連の手順を再現できるようにする。', 'size': 16, 'color': DARK, 'line': 1.25}])
# 右側：CSTの走行(通過域)を縦に
rx = Inches(7.1)
rect(s, rx, Inches(1.5), Inches(5.5), Inches(5.3), fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, rx, Inches(1.7), Inches(5.5), Inches(0.5),
        [{'text': 'CST の通過域（上 → 下）', 'size': 16, 'bold': True, 'color': NAVY, 'align': PP_ALIGN.CENTER}])
regions = ["M1 hand（一次運動野・手領域）", "Internal capsule（内包）",
           "Cerebral peduncle（大脳脚）", "Pons（橋）", "Medulla（延髄）"]
for i, r in enumerate(regions):
    y = Inches(2.35 + 0.78*i)
    bar = rect(s, rx + Inches(0.4), y, Inches(4.7), Inches(0.6),
               fill=CYAN if i % 2 == 0 else TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bar.adjustments[0] = 0.5
    textbox(s, rx + Inches(0.4), y, Inches(4.7), Inches(0.6),
            [{'text': r, 'size': 14, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 3. 全体ワークフロー
# ============================================================
s = slide(); bg(s)
header(s, None, "全体ワークフロー", "3")
steps = [
    ("1", "データ準備", "OSF から 5tt / wmfod_norm / T1_coreg を取得"),
    ("2", "mrview で表示", "T1 を開き overlay を重ねて確認"),
    ("3", "ROI を描く", "CST 通過域を塗り ROI Editor で保存"),
    ("4", "全脳トラクト生成", "tckgen で 100万本のストリームライン"),
    ("5", "ROI で絞り込み", "tckedit で CST のみ抽出"),
    ("6", "可視化", "全神経 → CST のみ を表示"),
]
cols = 3
cw, ch = Inches(3.9), Inches(2.35)
gx, gy = Inches(0.35), Inches(0.4)
x0, y0 = Inches(0.55), Inches(1.55)
for i, (n, t, d) in enumerate(steps):
    r, c = divmod(i, cols)
    x = x0 + c * (cw + gx)
    y = y0 + r * (ch + gy)
    card = rect(s, x, y, cw, ch, fill=LIGHT, line=PH_LN, line_w=Pt(0.75),
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.05
    num = rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.7), Inches(0.7),
               fill=BLUE, shape=MSO_SHAPE.OVAL)
    textbox(s, x + Inches(0.25), y + Inches(0.25), Inches(0.7), Inches(0.7),
            [{'text': n, 'size': 22, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER}],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x + Inches(1.1), y + Inches(0.28), cw - Inches(1.3), Inches(0.7),
            [{'text': t, 'size': 18, 'bold': True, 'color': NAVY}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x + Inches(0.3), y + Inches(1.15), cw - Inches(0.6), Inches(1.0),
            [{'text': d, 'size': 14, 'color': GRAY, 'line': 1.2}])

# ============================================================
# 共通：手順スライド生成
# ============================================================
def step_slide(step, title, idx, bullets, cmds=None, ph_label="スクリーンショット挿入予定",
               note=None):
    s = slide(); bg(s)
    header(s, step, title, idx)
    # 左：説明
    lx, lw = Inches(0.7), Inches(5.7)
    lines = []
    for b in bullets:
        lines.append({'text': '●  ' + b[0], 'size': 16, 'bold': True, 'color': BLUE,
                      'space_after': 2, 'line': 1.15})
        if len(b) > 1 and b[1]:
            lines.append({'text': '     ' + b[1], 'size': 14, 'color': DARK,
                          'space_after': 12, 'line': 1.2})
        else:
            lines[-1]['space_after'] = 12
    ly = Inches(1.45)
    textbox(s, lx, ly, lw, Inches(3.6), lines)
    if cmds:
        ch_in = 0.5 + 0.42*len(cmds)
        bottom_cap = 6.35 if note else 6.95
        cy_in = min(5.05, bottom_cap - ch_in)
        codebox(s, lx, Inches(cy_in), lw, Inches(ch_in), cmds)
    if note:
        textbox(s, lx, Inches(6.55), lw, Inches(0.7),
                [{'text': '📝 ' + note, 'size': 12, 'italic': True, 'color': GRAY, 'line': 1.15}])
    # 右：プレースホルダ
    placeholder(s, Inches(6.7), Inches(1.45), Inches(5.95), Inches(5.4), ph_label)
    return s

# 4. Step1 データ準備
step_slide("STEP 1", "データ準備（OSF Storage）", "4",
    [("必要なファイルを取得", "OSF Storage から解析に使う .mif を作業フォルダへ"),
     ("5tt_coreg.mif", "5組織タイプ分割（解剖学的拘束用）"),
     ("wmfod_norm.mif", "正規化済み白質 FOD（トラクト生成の方向場）"),
     ("T1_coreg.mif", "DWI と位置合わせ済みの T1 構造画像")],
    note="作業ディレクトリ: .../Tutrial/MRI/DWI/CST",
    ph_label="OSF Storage 画面のスクショ")

# 5. Step2 mrview
step_slide("STEP 2", "mrview で表示する", "5",
    [("作業フォルダへ移動", "ターミナルで対象ディレクトリへ cd"),
     ("mrview で T1 を開く", "構造画像を背景として表示"),
     ("overlay を重ねる", "Tool → Overlay から FOD などを重畳して確認")],
    cmds=["cd /Users/.../Tutrial/MRI/DWI/CST",
          "mrview T1_coreg.mif"],
    note="フルパスでファイルを指定して mrview を起動する。",
    ph_label="mrview / overlay のスクショ")

# 6. Step3 ROI
s = slide(); bg(s)
header(s, "STEP 3", "ROI を描く（ROI Editor）", "6")
textbox(s, Inches(0.7), Inches(1.45), Inches(5.7), Inches(3.8),
    [{'text': '●  CST の通過域を塗る', 'size': 16, 'bold': True, 'color': BLUE, 'space_after': 2},
     {'text': '     mrview の ROI Editor で各通過域をマスクとして描画', 'size': 14, 'color': DARK, 'space_after': 12, 'line': 1.2},
     {'text': '●  名前をつけて保存', 'size': 16, 'bold': True, 'color': BLUE, 'space_after': 2},
     {'text': '     領域ごとに <領域>_L.mif として書き出す', 'size': 14, 'color': DARK, 'space_after': 12, 'line': 1.2},
     {'text': '●  解剖の参照', 'size': 16, 'bold': True, 'color': BLUE, 'space_after': 2},
     {'text': '     Kandel『Principles of Neural Science』の図を参照して位置決め', 'size': 14, 'color': DARK, 'line': 1.2}])
# ROIテーブル
ty = Inches(4.55)
rois = [("Medulla_L", "延髄"), ("Pons_L", "橋"), ("Cerebral_peduncle_L", "大脳脚"),
        ("Internal_capsule_L", "内包"), ("M1_hand_L", "一次運動野(手)")]
rect(s, Inches(0.7), ty, Inches(5.7), Inches(0.45), fill=NAVY)
textbox(s, Inches(0.8), ty, Inches(3.0), Inches(0.45),
        [{'text': 'ROI ファイル名', 'size': 13, 'bold': True, 'color': WHITE}], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(4.0), ty, Inches(2.3), Inches(0.45),
        [{'text': '解剖部位', 'size': 13, 'bold': True, 'color': WHITE}], anchor=MSO_ANCHOR.MIDDLE)
for i, (f, jp) in enumerate(rois):
    ry = ty + Inches(0.45) + Inches(0.4)*i
    rect(s, Inches(0.7), ry, Inches(5.7), Inches(0.4),
         fill=WHITE if i % 2 else LIGHT, line=PH_LN, line_w=Pt(0.5))
    textbox(s, Inches(0.8), ry, Inches(3.1), Inches(0.4),
            [{'text': f, 'size': 12, 'bold': True, 'color': DARK, 'name': MONO}], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(4.0), ry, Inches(2.3), Inches(0.4),
            [{'text': jp, 'size': 12, 'color': GRAY}], anchor=MSO_ANCHOR.MIDDLE)
placeholder(s, Inches(6.7), Inches(1.45), Inches(5.95), Inches(5.4), "ROI Editor 操作画面のスクショ")

# ============================================================
# 7. 塗った領域：5つの ROI（実スクショ・ギャラリー）
# ============================================================
import os
ROI_DIR = os.path.join(os.path.dirname(__file__), '..', 'slides', 'roi_src')
# (ファイル名, ROI名, 日本語, マーカー色)  上 → 下の解剖順
roi_gallery = [
    ("roi_M1_hand_L.png",           "M1_hand_L",           "一次運動野・手領域", RGBColor(0x3B,0x6E,0xF0)),
    ("roi_Internal_capsule_L.png",  "Internal_capsule_L",  "内包",             RGBColor(0xE0,0x9B,0x3D)),
    ("roi_Cerebral_peduncle_L.png", "Cerebral_peduncle_L", "大脳脚",           RGBColor(0x2F,0xC2,0xD6)),
    ("roi_Pons_L.png",              "Pons_L",              "橋",               RGBColor(0x3B,0x6E,0xF0)),
    ("roi_Medulla_L.png",           "Medulla_L",           "延髄",             RGBColor(0x6A,0x55,0xD0)),
]
s = slide(); bg(s, RGBColor(0x0E, 0x16, 0x20))
rect(s, 0, 0, SW, Inches(1.15), fill=NAVY)
rect(s, 0, 0, Inches(0.18), SH, fill=CYAN)
textbox(s, Inches(0.55), 0, Inches(11.5), Inches(1.15),
        [{'text': '塗った領域：5つの ROI（CST の通過域）', 'size': 28, 'bold': True, 'color': WHITE}],
        anchor=MSO_ANCHOR.MIDDLE)
textbox(s, SW - Inches(1.2), Inches(0.32), Inches(0.9), Inches(0.5),
        [{'text': '7', 'size': 12, 'color': RGBColor(0xAF,0xC6,0xDB)}],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def roi_card(x, y, cw, ch, fname, name, jp, mc):
    card = rect(s, x, y, cw, ch, fill=RGBColor(0x00,0x00,0x00), line=mc, line_w=Pt(1.5),
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.04
    # ラベル帯
    lab_h = Inches(0.5)
    rect(s, x, y, cw, lab_h, fill=mc, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    dot = rect(s, x + Inches(0.14), y + Inches(0.15), Inches(0.2), Inches(0.2),
               fill=WHITE, shape=MSO_SHAPE.OVAL)
    textbox(s, x + Inches(0.42), y, cw - Inches(0.5), lab_h,
            [{'text': name + '   ' + jp, 'size': 12.5, 'bold': True, 'color': WHITE}],
            anchor=MSO_ANCHOR.MIDDLE)
    p = os.path.join(ROI_DIR, fname)
    if os.path.exists(p):
        add_image_fit(s, p, x + Inches(0.08), y + lab_h + Inches(0.05),
                      cw - Inches(0.16), ch - lab_h - Inches(0.13))

cw, ch = Inches(4.0), Inches(2.62)
gx = Inches(0.25)
# 上段3枚
y1 = Inches(1.42)
x0 = Inches(0.42)
for i in range(3):
    f, n, jp, mc = roi_gallery[i]
    roi_card(x0 + i*(cw+gx), y1, cw, ch, f, n, jp, mc)
# 下段2枚（中央寄せ）
y2 = Inches(4.22)
x0b = Emu(int((SW - (2*cw + gx)) / 2))
for i in range(2):
    f, n, jp, mc = roi_gallery[3+i]
    roi_card(x0b + i*(cw+gx), y2, cw, ch, f, n, jp, mc)
textbox(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.45),
        [{'text': '※ 各 ROI を mrview の ROI Editor で塗り、<領域>_L.mif として保存（解剖参照：Kandel）。'
                  '脳幹3領域（大脳脚・橋・延髄）の対応は要確認。',
          'size': 11, 'color': RGBColor(0x9D,0xB4,0xC8)}], anchor=MSO_ANCHOR.MIDDLE)

# 8. Step4 tckgen
step_slide("STEP 4", "全脳トラクト生成（tckgen）", "8",
    [("text editor でスクリプト作成", "tckgen_CST_L という名前でコマンドを記述"),
     ("シード", "gmwmSeed_coreg.mif（灰白質-白質境界）から発生"),
     ("ストリームライン数", "-select 1000000（100万本）を生成"),
     ("出力", "tracks_10bil.tck（全脳トラクト）")],
    cmds=["tckgen gmwmSeed_coreg.mif -select 1000000 \\",
          "       wmfod_norm.mif tracks_10bil.tck"],
    note="まずは全脳のストリームラインをまとめて生成しておく。",
    ph_label="text editor / コマンド記述のスクショ")

# 8. Step5 tckedit
step_slide("STEP 5", "ROI で絞り込む（tckedit）", "9",
    [("include で通過条件を指定", "5つの ROI を全て通る神経だけを残す"),
     ("入力", "tracks_10bil.tck（全脳トラクト）"),
     ("出力", "tracks_CST.tck（CST のみ）"),
     ("実行", "bash でスクリプトを指定して走らせる")],
    cmds=["tckedit -include Medulla_L.mif -include Pons_L.mif \\",
          "  -include Cerebral_peduncle_L.mif \\",
          "  -include Internal_capsule_L.mif -include M1_hand_L.mif \\",
          "  tracks_10bil.tck tracks_CST.tck",
          "bash .../CST/tckgen_CST_L.sh"],
    ph_label="ターミナル実行のスクショ")

# 9. Step6 可視化①
step_slide("STEP 6 ①", "可視化：全ての神経", "10",
    [("Tractography を選択", "mrview の Tool → Tractography を開く"),
     ("Run で得た2ファイルを読み込む", "生成したトラクトを表示"),
     ("結果", "脳全体の神経線維（全トラクト）が見える図")],
    note="まずは絞り込み前の全神経を確認する。",
    ph_label="全トラクト表示のスクショ")

# 10. Step6 可視化②
step_slide("STEP 6 ②", "可視化：CST のみ", "11",
    [("選んだ領域だけを通る神経を表示", "ROI を全て通過するトラクトのみにチェック"),
     ("結果", "選んだ通過域を通る神経 ＝ CST だけが表示される"),
     ("確認ポイント", "M1 から延髄まで連続して走行しているか")],
    note="tracks_CST.tck を読み込むと CST のみが描出される。",
    ph_label="CST のみ表示のスクショ")

# ============================================================
# 11. まとめ & 次のステップ（壁打ち）
# ============================================================
s = slide(); bg(s, NAVY)
rect(s, 0, 0, Inches(0.18), SH, fill=CYAN)
textbox(s, Inches(0.7), Inches(0.55), Inches(11), Inches(0.9),
        [{'text': 'まとめ & 次のステップ', 'size': 32, 'bold': True, 'color': WHITE}])
textbox(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(5.0),
        [{'text': '手順サマリ', 'size': 20, 'bold': True, 'color': CYAN, 'space_after': 10},
         {'text': '① データ準備 → ② mrview 表示 → ③ ROI 描画 →', 'size': 16, 'color': WHITE, 'space_after': 6, 'line': 1.3},
         {'text': '④ tckgen で全脳トラクト → ⑤ tckedit で CST 抽出 →', 'size': 16, 'color': WHITE, 'space_after': 6, 'line': 1.3},
         {'text': '⑥ Tractography で可視化（全神経 → CST のみ）', 'size': 16, 'color': WHITE, 'line': 1.3}])
box = rect(s, Inches(7.0), Inches(1.7), Inches(5.6), Inches(5.0), fill=RGBColor(0x1E,0x3A,0x5C),
           shape=MSO_SHAPE.ROUNDED_RECTANGLE)
box.adjustments[0] = 0.04
textbox(s, Inches(7.3), Inches(1.95), Inches(5.0), Inches(4.6),
        [{'text': '🛠 壁打ちポイント（要確認）', 'size': 18, 'bold': True, 'color': CYAN, 'space_after': 12},
         {'text': '・スクショ差し込み：別ファイルの全スクリーンショットを各手順へ配置（プレースホルダ枠を用意済み）', 'size': 14, 'color': WHITE, 'space_after': 10, 'line': 1.25},
         {'text': '・左右の別：今回は _L（左）。右側 CST も作るか？', 'size': 14, 'color': WHITE, 'space_after': 10, 'line': 1.25},
         {'text': '・コマンドの正確なオプション（mask/角度閾値等）を追記するか', 'size': 14, 'color': WHITE, 'space_after': 10, 'line': 1.25},
         {'text': '・対象者/発表時間に合わせて枚数・粒度を調整', 'size': 14, 'color': WHITE, 'line': 1.25}])

import os
out = os.path.join(os.path.dirname(__file__), '..', 'slides', 'CST_tractography_slides.pptx')
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved:", os.path.abspath(out), "slides:", len(prs.slides._sldIdLst))
